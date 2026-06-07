# -*- coding: utf-8 -*-
"""Construye el PPTX del proyecto de Programación en SIG (CGSM) desde los
fondos UNAL Verde. Espeja el deck web (14 diapositivas en español) e incrusta
el guion del orador en las notas de cada diapositiva."""
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from PIL import Image

# ---- palette (Monitor skin: teal mangrove + amber semaphore) ----
GREEN      = RGBColor(0x0F, 0x76, 0x6E)   # teal primary
GREEN_DEEP = RGBColor(0x0B, 0x5A, 0x54)
GOLD       = RGBColor(0xF5, 0x9E, 0x0B)   # amber accent
INK        = RGBColor(0x16, 0x24, 0x1D)
INK_SOFT   = RGBColor(0x5B, 0x6B, 0x64)
PAPER2     = RGBColor(0xF4, 0xF6, 0xF4)
LINE       = RGBColor(0xDD, 0xE1, 0xDC)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
WHITE_SOFT = RGBColor(0xDB, 0xE7, 0xDF)
TINT       = RGBColor(0xE7, 0xF1, 0xEB)
BLUE       = RGBColor(0x3B, 0x82, 0xC4)
PURPLE     = RGBColor(0x62, 0x0C, 0x78)

SERIF = "Georgia"
SANS  = "Calibri"

BRAND = r"C:/LINA/TOOLS/cult-dashboard-starter/public/brand"
FIG   = r"C:/LINA/TOOLS/cult-dashboard-starter/public/figures/prog"
OUT   = r"C:/LINA/UNAL 2626-1/PROG_SIG/Docker_1/proyecto-cgsm/docs/presentacion/Presentacion_ProgSIG_CGSM.pptx"
os.makedirs(os.path.dirname(OUT), exist_ok=True)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = 13.333, 7.5
BLANK = prs.slide_layouts[6]


def new_slide(bg):
    s = prs.slides.add_slide(BLANK)
    s.shapes.add_picture(os.path.join(BRAND, bg), 0, 0, Inches(SW), Inches(SH))
    return s


def notes(s, text):
    s.notes_slide.notes_text_frame.text = text


def box(s, x, y, w, h, anchor=None):
    tb = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    if anchor:
        tf.vertical_anchor = anchor
    return tf


BODY = 1.06
def _bs(size):
    return round(size * BODY * 2) / 2 if size <= 16 else size


def para(tf, text, size, color, bold=False, italic=False, font=SANS,
         align=PP_ALIGN.LEFT, after=4, first=False, line=None):
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    p.alignment = align
    if after is not None:
        p.space_after = Pt(after)
    p.space_before = Pt(0)
    if line is not None:
        p.line_spacing = line
    r = p.add_run()
    r.text = text
    f = r.font
    f.size = Pt(_bs(size))
    f.bold = bold
    f.italic = italic
    f.name = font
    f.color.rgb = color
    return p


def run(p, text, size, color, bold=False, italic=False, font=SANS):
    r = p.add_run()
    r.text = text
    f = r.font
    f.size = Pt(_bs(size))
    f.bold = bold
    f.italic = italic
    f.name = font
    f.color.rgb = color
    return r


def card(s, x, y, w, h, fill=PAPER2, line=LINE, line_w=1.0, radius=0.06):
    sh = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                            Inches(x), Inches(y), Inches(w), Inches(h))
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill
    if line is None:
        sh.line.fill.background()
    else:
        sh.line.color.rgb = line
        sh.line.width = Pt(line_w)
    sh.shadow.inherit = False
    try:
        sh.adjustments[0] = radius
    except Exception:
        pass
    return sh


def header(s, eyebrow, title, tsize=30):
    tf = box(s, 0.85, 0.62, 11.6, 0.35)
    para(tf, eyebrow.upper(), 11, GREEN, bold=True, font=SANS, after=0, first=True)
    tf2 = box(s, 0.85, 1.0, 11.6, 1.0)
    para(tf2, title, tsize, INK, bold=True, font=SERIF, after=0, first=True, line=1.04)


def bullets(s, x, y, w, h, items, size=12.5, color=INK_SOFT, gap=7,
            dot=GREEN, font=SANS):
    tf = box(s, x, y, w, h)
    for i, it in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(gap)
        p.space_before = Pt(0)
        p.line_spacing = 1.0
        r = p.add_run()
        r.text = "•  "
        r.font.size = Pt(_bs(size))
        r.font.name = font
        r.font.color.rgb = dot
        r.font.bold = True
        r2 = p.add_run()
        r2.text = it
        r2.font.size = Pt(_bs(size))
        r2.font.name = font
        r2.font.color.rgb = color


def stat(s, x, y, w, h, value, label, unit=None, vsize=23):
    card(s, x, y, w, h)
    tf = box(s, x + 0.18, y + 0.12, w - 0.34, h - 0.24, anchor=MSO_ANCHOR.MIDDLE)
    p = para(tf, value, vsize, GREEN, bold=True, font=SERIF, after=2, first=True)
    if unit:
        run(p, " " + unit, 11, INK_SOFT, bold=True, font=SANS)
    para(tf, label, 9.5, INK_SOFT, font=SANS, after=0, line=1.0)


def tag(s, x, y, w, h, text, fill=TINT, color=GREEN):
    card(s, x, y, w, h, fill=fill, line=None, radius=0.5)
    tf = box(s, x, y, w, h, anchor=MSO_ANCHOR.MIDDLE)
    para(tf, text.upper(), 9, color, bold=True, font=SANS, after=0,
         first=True, align=PP_ALIGN.CENTER)


def feature(s, x, y, w, h, title, body, accent=GREEN):
    card(s, x, y, w, h)
    bar = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                             Inches(x), Inches(y), Inches(0.08), Inches(h))
    bar.fill.solid(); bar.fill.fore_color.rgb = accent
    bar.line.fill.background(); bar.shadow.inherit = False
    tf = box(s, x + 0.28, y + 0.14, w - 0.46, h - 0.28, anchor=MSO_ANCHOR.MIDDLE)
    para(tf, title, 12.5, INK, bold=True, font=SANS, after=3, first=True, line=1.0)
    para(tf, body, 10.5, INK_SOFT, font=SANS, after=0, line=1.04)


def numbered(s, x, y, w, h, n, title, body, tagtxt=None):
    card(s, x, y, w, h)
    tf = box(s, x + 0.24, y + 0.16, w - 0.48, 0.5)
    para(tf, str(n), 20, GREEN, bold=True, font=SERIF, after=0, first=True)
    if tagtxt:
        tag(s, x + w - 1.25, y + 0.2, 1.05, 0.3, tagtxt)
    tf2 = box(s, x + 0.24, y + 0.74, w - 0.48, h - 0.9)
    para(tf2, title, 12.5, INK, bold=True, font=SANS, after=3, first=True, line=1.0)
    para(tf2, body, 10.5, INK_SOFT, font=SANS, after=0, line=1.04)


def fit_pic(s, path, x, y, bw, bh):
    """Centra y ajusta la imagen dentro del recuadro (x,y,bw,bh)."""
    iw, ih = Image.open(path).size
    ar = iw / ih
    box_ar = bw / bh
    if ar > box_ar:
        w = bw; h = bw / ar
    else:
        h = bh; w = bh * ar
    px = x + (bw - w) / 2
    py = y + (bh - h) / 2
    s.shapes.add_picture(path, Inches(px), Inches(py), Inches(w), Inches(h))


def fig(name):
    return os.path.join(FIG, name)


# ============================================================ 01 PORTADA
s = new_slide("monitor-title.png")
tf = box(s, 0.95, 1.25, 11.4, 0.4)
para(tf, "UNIVERSIDAD NACIONAL DE COLOMBIA", 12, GOLD, bold=True, after=3, first=True)
tf = box(s, 0.95, 1.72, 11.4, 0.35)
para(tf, "PROGRAMACIÓN EN SIG  ·  PROYECTO FINAL", 11.5, WHITE_SOFT, bold=True, after=0, first=True)
tf = box(s, 0.95, 2.35, 11.2, 2.1)
para(tf, "Pipeline multilenguaje para el monitoreo del manglar en la Ciénaga Grande de Santa Marta",
     33, WHITE, bold=True, font=SERIF, after=0, first=True, line=1.04)
gold = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.95), Inches(4.55), Inches(2.0), Inches(0.06))
gold.fill.solid(); gold.fill.fore_color.rgb = GOLD; gold.line.fill.background(); gold.shadow.inherit = False
tf = box(s, 0.95, 4.78, 8.4, 0.8)
para(tf, "Python · R · Julia sobre Google Earth Engine y Docker para caracterizar la dinámica espaciotemporal del manglar, 2013–2025.",
     14, WHITE_SOFT, font=SANS, after=0, first=True, line=1.12)
tf = box(s, 0.95, 5.85, 11.0, 1.2)
para(tf, "Lina María Quintero Fonseca", 17, WHITE, bold=True, font=SERIF, after=3, first=True)
para(tf, "Maestría en Geomática  ·  Facultad de Ciencias Agrarias · Sede Bogotá", 11.5, WHITE_SOFT, after=2)
para(tf, "Prof. Alexys H. Rodríguez-Avellaneda   ·   github.com/linaq11/proyecto-cgsm-curso", 11.5, WHITE_SOFT, after=0)
notes(s, "Buenas. Presento el proyecto final del curso de Programación en SIG: un pipeline multilenguaje, "
       "desarrollado en Python, R y Julia, para monitorear la dinámica del manglar de la Ciénaga Grande de "
       "Santa Marta entre 2013 y 2025. La idea central del proyecto es de programación: cómo orquestar tres "
       "lenguajes, una plataforma en la nube como Earth Engine y un contenedor Docker en un solo flujo "
       "reproducible. La Ciénaga es el caso de aplicación y también la línea de avance de mi tesis.")

# ============================================================ 02 CONTEXTO
s = new_slide("monitor-content.png")
header(s, "Contexto · Por qué importa", "La Ciénaga Grande de Santa Marta")
tf = box(s, 0.85, 2.15, 5.6, 1.5)
para(tf, "El complejo lagunar costero más extenso de Colombia, sitio Ramsar y Reserva de Biosfera, con ciclos "
     "documentados de degradación y recuperación del manglar ligados a la hipersalinización y al ENSO.",
     13.5, INK_SOFT, font=SANS, after=0, first=True, line=1.18)
stat(s, 0.85, 3.95, 2.75, 1.05, "835,3", "AOI acotado (SFF CGSM + VPI Salamanca)", unit="km²")
stat(s, 3.72, 3.95, 2.75, 1.05, "2013–2025", "Serie multisensor analizada", vsize=19)
stat(s, 0.85, 5.12, 2.75, 1.05, "Ramsar", "Reserva de Biosfera UNESCO (2000)", unit="1998")
stat(s, 3.72, 5.12, 2.75, 1.05, "ENSO", "El Niño 2015–16 · La Niña 2020–22", vsize=20)
card(s, 6.85, 2.15, 5.6, 4.0, fill=PAPER2)
tf = box(s, 7.1, 2.4, 5.1, 0.4)
para(tf, "El reto de monitoreo", 14, INK, bold=True, after=0, first=True)
bullets(s, 7.1, 3.0, 5.15, 2.2, [
    "Hay datos satelitales frecuentes, pero no se convierten en monitoreo espacialmente explícito y reproducible.",
    "La nubosidad tropical persistente limita la observación óptica.",
    "El radar aporta continuidad, pero exige integrar varias fuentes y lenguajes.",
], size=11.5, gap=9)
card(s, 7.1, 5.25, 5.1, 0.78, fill=TINT, line=None)
tf = box(s, 7.32, 5.4, 4.66, 0.5, anchor=MSO_ANCHOR.MIDDLE)
p = para(tf, "La mortandad de manglar de septiembre de 2020 ", 11, INK, bold=True, after=0, first=True, line=1.04)
run(p, "muestra la necesidad de un seguimiento espacial y continuo.", 11, INK_SOFT)
notes(s, "La Ciénaga Grande de Santa Marta es el complejo lagunar costero más extenso de Colombia, sitio Ramsar "
       "y Reserva de Biosfera. Su manglar ha sufrido ciclos de degradación y recuperación asociados a la "
       "hipersalinización y a la variabilidad del ciclo El Niño-Oscilación del Sur. El problema no es la falta "
       "de datos satelitales, que son frecuentes, sino que esa información no se convierte en un monitoreo "
       "espacialmente explícito, reproducible y que integre lo óptico, lo radar y lo climático en una sola "
       "arquitectura.")

# ============================================================ 03 VACÍO / PREGUNTA
s = new_slide("monitor-content.png")
header(s, "Justificación · La oportunidad", "El vacío no es de datos, es de integración")
xs = 0.85
data3 = [
    ("1", "Observación fragmentada", "Lo óptico y el radar se analizan por separado; falta articularlos en una serie consistente.", "Óptico + radar"),
    ("2", "Sin flujo reproducible", "Segmentación, fragmentación y quiebres temporales no viven en una misma arquitectura ejecutable.", "Pipeline"),
    ("3", "Forzamientos sin acoplar", "Las perturbaciones rara vez se cruzan con ENSO, caudal y precipitación en el mismo marco.", "Clima"),
]
for n, t, d, tg in data3:
    numbered(s, xs, 2.15, 3.78, 2.35, n, t, d, tagtxt=tg)
    xs += 3.95
card(s, 0.85, 4.75, 11.6, 1.55, fill=GREEN, line=None)
tf = box(s, 1.2, 4.92, 11.0, 0.32)
para(tf, "PREGUNTA DE INVESTIGACIÓN", 10.5, GOLD, bold=True, after=0, first=True)
tf = box(s, 1.2, 5.28, 11.0, 0.95, anchor=MSO_ANCHOR.MIDDLE)
para(tf, "¿Cómo varió la cobertura, la fragmentación y el vigor del manglar de la CGSM entre 2013 y 2025, y cómo se "
     "relacionan las perturbaciones detectadas con los forzamientos climáticos asociados al ENSO?",
     15.5, WHITE, font=SERIF, after=0, first=True, line=1.2)
notes(s, "El vacío que abordo no es la ausencia de información sobre el manglar de la Ciénaga, sino la falta de un "
       "flujo reproducible que articule, en una misma arquitectura, observación óptica y radar, segmentación "
       "automática, métricas de fragmentación, detección de quiebres temporales y forzamientos climático-"
       "hidrológicos. De ahí la pregunta de investigación: cómo varió la cobertura, la fragmentación y el vigor "
       "del manglar entre 2013 y 2025, y cómo se relacionan las perturbaciones con los forzamientos climáticos "
       "del periodo.")

# ============================================================ 04 OBJETIVOS
s = new_slide("monitor-content.png")
header(s, "Objetivos · Una meta, cinco pasos", "Qué construye el pipeline")
card(s, 0.85, 2.15, 4.7, 4.15, fill=GREEN, line=None)
tf = box(s, 1.12, 2.45, 4.18, 0.32)
para(tf, "OBJETIVO GENERAL", 10.5, GOLD, bold=True, after=0, first=True)
tf = box(s, 1.12, 2.95, 4.18, 2.4)
para(tf, "Desarrollar un pipeline multilenguaje para el monitoreo del manglar de la Ciénaga Grande de Santa Marta "
     "entre 2013 y 2025.", 16, WHITE, font=SERIF, after=8, first=True, line=1.28)
para(tf, "Reproducible, interoperable y transferible a otros humedales costeros tropicales.",
     11, WHITE_SOFT, font=SANS, after=0, line=1.1)
objs = [
    ("1 · Datacube multitemporal", "Índices NDVI, NDWI y CMRI sobre la CGSM 2013–2025 (Landsat 8 + Sentinel-2)."),
    ("2 · Degradación y recuperación", "Anomalías z-score y quiebres estructurales BFAST sobre las series de NDVI."),
    ("3 · Segmentación y fragmentación", "SamGeo y métricas del paisaje calculadas en EPSG:9377."),
    ("4 · Forzamientos climático-hidrológicos", "Relación con ERA5-Land, ENSO, caudal IDEAM y precipitación CHIRPS."),
    ("5 · Validación cartográfica", "Umbrales vs. Random Forest frente a INVEMAR 1:25.000 y ESA WorldCover."),
]
yy = 2.15
for t, d in objs:
    feature(s, 5.85, yy, 6.6, 0.76, t, d)
    yy += 0.84
notes(s, "El objetivo general es desarrollar ese pipeline multilenguaje de monitoreo. Se concreta en cinco objetivos "
       "específicos: construir un datacube multitemporal de índices espectrales; detectar degradación y recuperación "
       "con anomalías z y quiebres BFAST; segmentar el manglar con SamGeo y medir su fragmentación en EPSG 9377; "
       "evaluar la relación con forzamientos climático-hidrológicos; y validar la clasificación frente a INVEMAR y "
       "ESA WorldCover, comparando umbrales con un Random Forest.")

# ============================================================ 05 ÁREA DE ESTUDIO
s = new_slide("monitor-content.png")
header(s, "Área de estudio · AOI acotado", "835,3 km² de área protegida oficial")
feature(s, 0.85, 2.2, 5.7, 1.18, "AOI acotado RUNAP",
        "Unión del Santuario de Fauna y Flora CGSM (26.810 ha) y el Vía Parque Isla de Salamanca. Reemplaza un AOI preliminar de 5.073 km².")
feature(s, 0.85, 3.5, 5.7, 1.18, "8 estaciones de muestreo",
        "Cinco INVEMAR-GBIF (limnológicas) y tres complementarias sobre manglar verificado por NDVI > 0,4.")
feature(s, 0.85, 4.8, 5.7, 1.18, "Áreas en EPSG:9377",
        "MAGNA-SIRGAS, sistema oficial colombiano. ~51% del AOI es agua permanente o estacional (JRC).")
card(s, 6.75, 2.05, 5.75, 4.55, fill=WHITE)
fit_pic(s, fig("area.png"), 6.95, 2.2, 5.35, 4.25)
notes(s, "El área de estudio se delimitó desde el inicio para no quedar abstracta. Trabajo sobre el AOI acotado de "
       "835,3 kilómetros cuadrados, que corresponde a la unión del Santuario de Fauna y Flora y el Vía Parque Isla "
       "de Salamanca, tomados del RUNAP. Un AOI preliminar de cinco mil kilómetros inflaba el manglar potencial con "
       "vegetación riberana y salitrales, así que todos los resultados se sostienen en el AOI acotado. Para las "
       "series definí ocho estaciones: cinco del INVEMAR publicadas en GBIF y tres complementarias sobre manglar "
       "verificado. Las áreas se calculan en EPSG 9377, el sistema oficial colombiano.")

# ============================================================ 06 FUENTES DE DATOS
s = new_slide("monitor-content.png")
header(s, "Datos · 13 conjuntos, 6 categorías", "Una base multisensor y multifuente")
src = [
    ("Óptico satelital", "Sentinel-2 MSI (789 img, 10 m) y Landsat 8/9 (serie NDVI 2013–2017)."),
    ("Radar Sentinel-1", "SAR banda VH, 10 m. Inundación bajo dosel por doble rebote."),
    ("Cartografía de referencia", "ESA WorldCover v200 e INVEMAR 1:25.000 para validación."),
    ("Clima e hidrología", "ERA5-Land, ENSO ONI/SOI (NOAA), caudal IDEAM, CHIRPS."),
    ("Agua e inundación", "JRC Global Surface Water y Global Flood Database (16 eventos)."),
    ("Elevación y campo", "SRTM v3 (<10 m) e INVEMAR-GBIF (estructura forestal)."),
]
xs = [0.85, 4.72, 8.59]
for i, (t, d) in enumerate(src):
    col = i % 3
    row = i // 3
    feature(s, xs[col], 2.2 + row * 1.5, 3.62, 1.32, t, d)
card(s, 0.85, 5.35, 11.6, 0.95, fill=PAPER2)
tf = box(s, 1.1, 5.5, 7.6, 0.66, anchor=MSO_ANCHOR.MIDDLE)
p = para(tf, "La mayoría se adquiere desde Google Earth Engine ", 11.5, GREEN, bold=True, after=0, first=True, line=1.05)
run(p, "(sin descargar imágenes pesadas).", 11.5, INK_SOFT)
tag(s, 9.0, 5.62, 1.05, 0.42, "EPSG:9377")
tag(s, 10.15, 5.62, 1.15, 0.42, "NetCDF CF-1.8")
tag(s, 11.4, 5.62, 1.0, 0.42, "10-30 m")
notes(s, "El pipeline integra trece fuentes en seis categorías: óptico Sentinel-2 y Landsat 8, radar Sentinel-1, "
       "cartografías de referencia como ESA WorldCover y la oficial de INVEMAR, el modelo de elevación SRTM, datos "
       "de campo del INVEMAR, e índices climático-hidrológicos: ERA5-Land, ENSO de la NOAA, caudal del IDEAM y "
       "precipitación CHIRPS. La mayoría se obtuvo desde Google Earth Engine, lo que evita descargar imágenes "
       "pesadas. La validación se hace contra dos cartografías independientes para tener un techo metodológico "
       "realista.")

# ============================================================ 07 PIPELINE MULTILENGUAJE
s = new_slide("monitor-content.png")
header(s, "Arquitectura · El corazón del proyecto", "5 fuentes → datacube → 3 lenguajes → 5 productos")
card(s, 0.85, 2.15, 7.05, 4.15, fill=WHITE)
fit_pic(s, fig("arquitectura.png"), 1.0, 2.3, 6.75, 3.85)
feature(s, 8.15, 2.15, 4.3, 1.02, "Python · adquisición y modelado",
        "GEE, SamGeo, Random Forest, datacubes xarray y dashboard.", accent=BLUE)
feature(s, 8.15, 3.27, 4.3, 1.02, "R · series y estadística",
        "Quiebres BFAST, cubo stars y réplicas de ENSO y caudal.", accent=PURPLE)
feature(s, 8.15, 4.39, 4.3, 1.02, "Julia · geometría y topología",
        "Fragmentación, áreas EPSG:9377 y predicados DE-9IM (LibGEOS.jl).", accent=GOLD)
card(s, 8.15, 5.52, 4.3, 0.78, fill=GREEN, line=None)
tf = box(s, 8.4, 5.62, 3.85, 0.6, anchor=MSO_ANCHOR.MIDDLE)
p = para(tf, "Encuentro: datacube NetCDF CF-1.8. ", 10.5, GOLD, bold=True, after=0, first=True, line=1.05)
run(p, "Todo en el contenedor Docker sig_unal v1.11.", 10.5, WHITE)
notes(s, "Esta es la diapositiva central del proyecto de programación. El flujo asigna a cada lenguaje su fortaleza: "
       "Python centraliza la adquisición en Earth Engine, la segmentación con SamGeo, el Random Forest, los "
       "datacubes y el dashboard; R hace la detección de quiebres con BFAST, el cubo stars y las réplicas "
       "estadísticas de ENSO y caudal; y Julia se encarga del cómputo geométrico, las métricas de fragmentación y "
       "los predicados topológicos DE-9IM con LibGEOS. El punto de encuentro es un datacube en NetCDF CF-1.8, "
       "reproyectado a EPSG 9377, que los tres lenguajes leen sin volver a consultar la API. Todo corre dentro del "
       "contenedor Docker sig_unal v1.11.")

# ============================================================ 08 MÉTODOS
s = new_slide("monitor-content.png")
header(s, "Métodos · 5 fases + módulo climático", "Decisiones técnicas del flujo")
meth = [
    ("1", "Datacube", "Composiciones Sentinel-2 (nubes <20%) y Landsat 8/9; NetCDF CF-1.8 a 30 m, EPSG:9377.", "Fase 1"),
    ("2", "Series + BFAST", "929 obs. de NDVI; z-score (z<-2) y quiebres BFAST con h = 0,15 y 0,10.", "Fase 2"),
    ("3", "Segmentación", "SamGeo backbone vit_b sobre RGB de 3 periodos; filtro de área 1–5.000 ha.", "Fase 3"),
    ("4", "Fragmentación", "Julia: nº de parches, MSI, NND y predicados DE-9IM (intersects, contains).", "Fase 4"),
    ("S1", "Inundación SAR", "Sentinel-1 VH: SAR diff = seco − inundado; doble rebote bajo dosel.", "Módulo"),
    ("5", "Random Forest", "100 árboles, 15 variables, 1.000 puntos, validación cruzada K = 5.", "Fase 5"),
]
xs = [0.85, 4.72, 8.59]
for i, (n, t, d, tg) in enumerate(meth):
    col = i % 3
    row = i // 3
    numbered(s, xs[col], 2.2 + row * 2.1, 3.62, 1.92, n, t, d, tagtxt=tg)
notes(s, "En cuanto a los métodos, el flujo se organiza en cinco fases más un módulo climático transversal. Construyo "
       "el datacube con composiciones Sentinel-2 y Landsat; extraigo series mensuales de NDVI y aplico z-score y "
       "BFAST con ventanas de doce y dieciocho meses; segmento con SamGeo usando el backbone vit_b; calculo "
       "fragmentación en Julia con número de parches, índice de forma y distancia al vecino; detecto inundación con "
       "Sentinel-1 comparando el backscatter seco contra el inundado; y entreno un Random Forest de cien árboles con "
       "quince variables como referencia frente a la regla por umbrales.")

# ============================================================ 09 RESULTADOS 1 NDVI/BFAST
s = new_slide("monitor-content.png")
header(s, "Resultados 1 · Dinámica temporal", "NDVI, anomalías y quiebres BFAST")
card(s, 0.85, 2.15, 7.0, 4.15, fill=WHITE)
fit_pic(s, fig("ndvi_serie.png"), 1.0, 2.55, 6.7, 3.35)
stat(s, 8.1, 2.15, 2.12, 1.02, "929", "Obs. mensuales de NDVI")
stat(s, 10.35, 2.15, 2.12, 1.02, "18", "Anomalías (z<-2)")
stat(s, 8.1, 3.3, 2.12, 1.02, "0,80→0,60→0,80", "NDVI mediano manglar", vsize=11.5)
stat(s, 10.35, 3.3, 2.12, 1.02, "2016", "Quiebre BFAST (El Niño)")
card(s, 8.1, 4.45, 4.37, 1.85, fill=PAPER2)
tf = box(s, 8.32, 4.6, 3.95, 0.34)
para(tf, "Lectura clave", 12.5, INK, bold=True, after=0, first=True)
tf = box(s, 8.32, 5.0, 3.95, 1.2)
para(tf, "La caída 2019–2020 coincide con la sequía y La Niña. Restringir BFAST a las cuatro estaciones de manglar "
     "denso afina la señal de septiembre 2020 y revela respuestas en 2022 y 2023–24.",
     10.5, INK_SOFT, font=SANS, after=0, first=True, line=1.12)
notes(s, "Los primeros resultados son temporales. La serie combinó 929 observaciones mensuales de NDVI y reveló "
       "dieciocho anomalías significativas, concentradas en septiembre de 2020, marzo de 2018 y 2016. El NDVI "
       "mediano del manglar osciló en torno a 0,80, cayó a 0,60 entre 2019 y 2020 por la sequía y el inicio de La "
       "Niña, y se estabilizó de nuevo en 0,80 desde 2022. BFAST detectó un quiebre generalizado en 2016, asociado "
       "a El Niño. Al restringir el análisis a las estaciones de manglar denso se afinó la señal de 2020 y se vieron "
       "respuestas posteriores coherentes con fases ENSO contrastantes.")

# ============================================================ 10 RESULTADOS 2 FRAGMENTACIÓN
s = new_slide("monitor-content.png")
header(s, "Resultados 2 · Dinámica espacial", "Contracción con consolidación estructural")
card(s, 0.85, 2.15, 5.5, 4.15, fill=WHITE)
fit_pic(s, fig("ndvi_cambio.png"), 1.0, 2.3, 5.2, 3.85)
stat(s, 6.6, 2.15, 2.85, 1.02, "12.426→4.037", "Cobertura clasificada", unit="ha", vsize=16)
stat(s, 9.62, 2.15, 2.85, 1.02, "79→15", "Número de parches")
stat(s, 6.6, 3.3, 2.85, 1.02, "157→269", "Área media por parche", unit="ha", vsize=18)
stat(s, 9.62, 3.3, 2.85, 1.02, "+76", "Cambio neto", unit="km²")
card(s, 6.6, 4.45, 5.87, 1.85, fill=PAPER2)
tf = box(s, 6.85, 4.6, 5.4, 0.34)
para(tf, "No es pérdida directa de cobertura", 12.5, INK, bold=True, after=0, first=True)
tf = box(s, 6.85, 5.0, 5.4, 1.2)
para(tf, "Menos parches pero más densos: MSI 0,51→1,46 y NND 1,10→2,39 km. La contracción del área clasificada "
     "se lee junto con la ganancia de vigor del NDVI. Topología DE-9IM idéntica en Python y Julia.",
     10.5, INK_SOFT, font=SANS, after=0, first=True, line=1.12)
notes(s, "En lo espacial, la segmentación con SamGeo arrojó un cambio neto positivo de 76 kilómetros cuadrados entre "
       "la degradación y el estado actual. Las métricas de fragmentación muestran algo interesante: la cobertura "
       "clasificada cayó de 12.426 a 4.037 hectáreas y el número de parches bajó de 79 a 15, pero el área media por "
       "parche subió de 157 a 269 hectáreas. Es decir, el manglar se reorganizó en menos parches pero más densos y "
       "maduros. La contracción del área clasificada no es pérdida directa de cobertura: hay que leerla junto con la "
       "ganancia de vigor del NDVI. La topología DE-9IM se validó de forma idéntica entre Python y Julia.")

# ============================================================ 11 RESULTADOS 3 CLIMA + SAR
s = new_slide("monitor-content.png")
header(s, "Resultados 3 · Forzamiento e inundación", "Acoplamiento climático y radar")
card(s, 0.85, 2.15, 6.7, 4.15, fill=WHITE)
fit_pic(s, fig("sar.png"), 1.0, 2.5, 6.4, 3.45)
stat(s, 7.8, 2.15, 2.27, 1.02, "ρ +0,256", "Caudal Magdalena ↔ NDVI (3 m)", vsize=18)
stat(s, 10.2, 2.15, 2.27, 1.02, "59,02", "Inundación SAR sep. 2020", unit="km²")
stat(s, 7.8, 3.3, 2.27, 1.02, "43,08", "Bajo dosel (doble rebote)", unit="km²")
stat(s, 10.2, 3.3, 2.27, 1.02, "ρ +0,81", "SAR-VH ↔ NDVI (Pajarales)", vsize=18)
card(s, 7.8, 4.45, 4.67, 1.85, fill=PAPER2)
tf = box(s, 8.02, 4.6, 4.25, 0.34)
para(tf, "Desagregar por naturaleza espectral", 12, INK, bold=True, after=0, first=True)
tf = box(s, 8.02, 5.0, 4.25, 1.2)
para(tf, "Manglar y cuerpos de agua responden con signos opuestos: promediar cancela la señal. El radar confirma la "
     "recuperación estructural del dosel que el NDVI satura.",
     10.5, INK_SOFT, font=SANS, after=0, first=True, line=1.12)
notes(s, "El forzamiento climático se evaluó con cuatro fuentes y desagregando las estaciones por su naturaleza "
       "espectral, porque promediar manglar y cuerpos de agua cancela señales opuestas. El aporte de agua dulce del "
       "río Magdalena se correlaciona positivamente con el vigor del manglar, con rho de 0,256 a tres meses, y "
       "CHIRPS lo confirma de forma independiente. En el radar, la detección con Sentinel-1 para septiembre de 2020 "
       "identificó 59 kilómetros cuadrados afectados, de los cuales 43 corresponden a inundación bajo dosel por "
       "doble rebote agua-tronco. Y la serie mensual de SAR muestra recuperación estructural del dosel, con "
       "correlaciones de hasta 0,8 con el NDVI en el Complejo de Pajarales.")

# ============================================================ 12 RESULTADOS 4 VALIDACIÓN / RF
s = new_slide("monitor-content.png")
header(s, "Resultados 4 · Validación cartográfica", "Umbrales vs. Random Forest")
card(s, 0.85, 2.15, 5.6, 4.15, fill=WHITE)
fit_pic(s, fig("rf.png"), 1.0, 2.35, 5.3, 3.75)
card(s, 6.7, 2.15, 5.77, 2.55, fill=PAPER2)
tf = box(s, 6.95, 2.32, 5.3, 0.34)
para(tf, "F1-score por método y referencia", 13, INK, bold=True, after=0, first=True)
rows = [
    ("Umbrales vs. INVEMAR", "0,583"),
    ("Umbrales vs. WorldCover", "0,548"),
    ("Random Forest vs. INVEMAR", "0,826"),
    ("Random Forest vs. WorldCover", "0,889"),
]
ry = 2.86
for k, v in rows:
    tf = box(s, 6.95, ry, 4.0, 0.36, anchor=MSO_ANCHOR.MIDDLE)
    para(tf, k, 11.5, INK_SOFT, font=SANS, after=0, first=True)
    tf = box(s, 11.0, ry, 1.25, 0.36, anchor=MSO_ANCHOR.MIDDLE)
    para(tf, v, 14, GREEN, bold=True, font=SERIF, after=0, first=True, align=PP_ALIGN.RIGHT)
    ry += 0.44
stat(s, 6.7, 4.95, 2.83, 1.35, "0,833", "Techo metodológico (INVEMAR ↔ WorldCover)")
stat(s, 9.64, 4.95, 2.83, 1.35, "SWIR + dist_agua", "Pesan más que el NDVI (Gini)", vsize=15)
notes(s, "Para la validación cartográfica comparo la clasificación contra dos referencias: INVEMAR a 1:25.000 y ESA "
       "WorldCover. La regla por umbrales es conservadora, con F1 de 0,58 y 0,55. Como las dos cartografías oficiales "
       "coinciden entre sí en un F1 de 0,83, ese es el techo metodológico realista, no el uno perfecto. El Random "
       "Forest, entrenado con mil puntos y validación cruzada, mejora de forma marcada: F1 de 0,83 frente a INVEMAR "
       "y 0,89 frente a WorldCover, sobre todo duplicando la sensibilidad. La importancia de variables muestra que "
       "las bandas SWIR y la distancia al agua pesan tanto o más que el NDVI, coherente con un manglar inundado y "
       "halófito.")

# ============================================================ 13 ALERTAS + REPRODUCIBILIDAD
s = new_slide("monitor-content.png")
header(s, "Operación · Hacia un Gemelo Digital", "Alertas tempranas y reproducibilidad")
card(s, 0.85, 2.15, 6.6, 4.15, fill=WHITE)
fit_pic(s, fig("semaforo.png"), 1.0, 2.3, 6.3, 3.85)
stat(s, 7.7, 2.15, 1.55, 1.1, "5", "Estables")
stat(s, 9.4, 2.15, 1.5, 1.1, "3", "En alerta")
stat(s, 11.0, 2.15, 1.47, 1.1, "0", "Críticas")
feature(s, 7.7, 3.45, 4.77, 1.32, "Validación trilingüe",
        "Python, R y Julia producen valores idénticos hasta 3 decimales (rho y topología DE-9IM).")
feature(s, 7.7, 4.92, 4.77, 1.38, "Reproducible de extremo a extremo",
        "Dashboard HTML de 17 capas, contenedor Docker sig_unal v1.11 y repositorio público en GitHub.")
notes(s, "El pipeline cierra con un módulo de alertas tempranas, que es el componente operacional del futuro Gemelo "
       "Digital. Usa bfastmonitor en modo incremental sobre las series y clasifica cada estación con una lógica de "
       "semáforo. Al cierre de 2025 hay cinco estaciones estables, tres en alerta y ninguna crítica, con el núcleo "
       "de manglar denso del Complejo de Pajarales en verde. Y todo el proyecto es reproducible: la validación "
       "trilingüe Python, R y Julia produce valores idénticos hasta tres decimales, el contenedor Docker congela las "
       "versiones, y un dashboard HTML de diecisiete capas resume los productos.")

# ============================================================ 14 MONITOR EN VIVO
s = new_slide("monitor-content.png")
header(s, "Producto · Monitor en vivo", "El pipeline termina en un monitor consultable")
card(s, 0.85, 2.15, 7.05, 4.15, fill=WHITE)
fit_pic(s, fig("monitor.png"), 0.98, 2.28, 6.79, 3.89)
card(s, 8.15, 2.15, 4.3, 2.62, fill=PAPER2)
tf = box(s, 8.4, 2.32, 3.85, 0.34)
para(tf, "Cinco pestañas", 13, INK, bold=True, after=0, first=True)
bullets(s, 8.4, 2.8, 3.9, 1.9, [
    "Resumen — estado y métricas clave",
    "Cobertura — NDVI, cambio y mapa interactivo",
    "Clima e hidrología — ENSO, caudal, CHIRPS",
    "Validación multilenguaje — Python · R · Julia",
    "Acerca de — datos, métodos y repositorio",
], size=10.5, gap=6)
stat(s, 8.15, 4.95, 2.07, 0.95, "17", "Capas temáticas sobre el AOI", vsize=22)
stat(s, 10.38, 4.95, 2.07, 0.95, "HTML", "Autocontenido, sin Jupyter", vsize=18)
card(s, 8.15, 6.05, 4.3, 0.55, fill=GREEN, line=None)
tf = box(s, 8.4, 6.11, 3.95, 0.44, anchor=MSO_ANCHOR.MIDDLE)
p = para(tf, "↗  ", 12, GOLD, bold=True, after=0, first=True)
run(p, "linaq11.github.io/proyecto-cgsm-curso/dashboard.html", 9.5, WHITE, bold=True)
notes(s, "Todo el pipeline termina en un producto consultable: el CGSM Monitor, un dashboard HTML publicado en "
       "GitHub Pages. Tiene cinco pestañas: Resumen, Cobertura, Clima e hidrología, Validación multilenguaje y "
       "Acerca de. Reúne los indicadores clave del estado del manglar al cierre de 2025: 4.037 hectáreas de manglar "
       "protegido, 835 kilómetros cuadrados de área, 59 inundados en 2020, el vínculo entre caudal y manglar y trece "
       "años de monitoreo continuo. Integra diecisiete capas temáticas y cualquiera puede navegarlo en el navegador, "
       "sin Jupyter ni software especializado. Es la cara operativa del flujo y el primer paso hacia un Gemelo "
       "Digital de la Ciénaga.")

# ============================================================ 15 CONCLUSIONES
s = new_slide("monitor-title.png")
tf = box(s, 0.95, 1.35, 11.4, 0.4)
para(tf, "CONCLUSIONES", 12, GOLD, bold=True, after=0, first=True)
tf = box(s, 0.95, 1.95, 11.4, 1.5)
para(tf, "Un flujo reproducible que orquesta tres lenguajes para monitorear el manglar",
     26, WHITE, bold=True, font=SERIF, after=0, first=True, line=1.1)
gold = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.95), Inches(3.5), Inches(1.8), Inches(0.06))
gold.fill.solid(); gold.fill.fore_color.rgb = GOLD; gold.line.fill.background(); gold.shadow.inherit = False
concl = [
    ("Hallazgos", "Respuesta a El Niño 2016 y La Niña 2020; contracción con consolidación; RF mejora la regla por umbrales."),
    ("Límites", "Sin campo independiente; red hidrométrica limitada; relaciones de asociación temporal, no causalidad."),
    ("Trabajo futuro", "Validación in situ, más estaciones y extensión hacia un Gemelo Digital costero con alertas."),
]
xc = 0.95
for t, d in concl:
    sh = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(xc), Inches(3.95), Inches(3.78), Inches(1.95))
    sh.fill.solid(); sh.fill.fore_color.rgb = GREEN_DEEP; sh.line.color.rgb = GOLD; sh.line.width = Pt(0.75)
    sh.shadow.inherit = False
    try: sh.adjustments[0] = 0.06
    except Exception: pass
    tf = box(s, xc + 0.28, 4.18, 3.3, 0.4)
    para(tf, t, 14, GOLD, bold=True, font=SANS, after=0, first=True)
    tf = box(s, xc + 0.28, 4.66, 3.3, 1.15)
    para(tf, d, 11, WHITE_SOFT, font=SANS, after=0, first=True, line=1.16)
    xc += 3.95
tf = box(s, 0.95, 6.35, 11.4, 0.5)
para(tf, "Lina María Quintero Fonseca   ·   Monitor: linaq11.github.io/proyecto-cgsm-curso/dashboard.html   ·   Gracias.",
     12, WHITE_SOFT, font=SANS, after=0, first=True)
notes(s, "Para cerrar: el pipeline demuestra que se puede orquestar Python, R y Julia en un flujo reproducible que "
       "articula evidencia espectral, espacial e hidroclimática del manglar. Los hallazgos clave son la respuesta "
       "del manglar a El Niño 2016 y La Niña 2020, una contracción con consolidación estructural, y la mejora del "
       "Random Forest sobre la regla por umbrales. Reconozco limitaciones: no hubo campo, la red hidrométrica es "
       "limitada, y las relaciones son asociación temporal, no causalidad. Como trabajo futuro, validación in situ, "
       "más estaciones y extender la ventana hacia un Gemelo Digital costero. Gracias, quedo atenta a sus preguntas.")

prs.save(OUT)
print("SAVED", OUT, "slides:", len(prs.slides._sldIdLst))
