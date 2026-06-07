# -*- coding: utf-8 -*-
"""Build the thesis-defense PPTX from the official UNAL Verde backgrounds.
Mirrors the web deck (13 slides) and embeds the Spanish speaker script in
each slide's notes."""
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ---- palette (UNAL Verde) ----
GREEN      = RGBColor(0x0B, 0x7A, 0x4B)
GREEN_DEEP = RGBColor(0x06, 0x53, 0x2F)
GOLD       = RGBColor(0xFF, 0xB9, 0x3E)
INK        = RGBColor(0x29, 0x27, 0x2D)
INK_SOFT   = RGBColor(0x5B, 0x55, 0x60)
PAPER2     = RGBColor(0xF4, 0xF6, 0xF4)
LINE       = RGBColor(0xDD, 0xE1, 0xDC)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
WHITE_SOFT = RGBColor(0xDB, 0xE7, 0xDF)
TINT       = RGBColor(0xE7, 0xF1, 0xEB)
GOLD_INK   = RGBColor(0x8A, 0x5E, 0x00)

SERIF = "Georgia"
SANS  = "Calibri"

BRAND = r"C:/LINA/TOOLS/cult-dashboard-starter/public/brand"
FIG   = r"C:/LINA/TOOLS/cult-dashboard-starter/public/figures"
OUT   = r"C:/LINA/UNAL 2626-1/AVANCE TESIS/Defensa_Lina_GeoAI_CGSM.pptx"

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = 13.333, 7.5
BLANK = prs.slide_layouts[6]


def new_slide(bg):
    s = prs.slides.add_slide(BLANK)
    s.shapes.add_picture(os.path.join(BRAND, bg), 0, 0, Inches(SW), Inches(SH))
    return s


def notes(s, text):
    s.notes_slide.notes_text_frame.text = text


def box(s, x, y, w, h, anchor=None):
    tb = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    if anchor:
        tf.vertical_anchor = anchor
    return tf


BODY = 1.07  # body-text scale; display text (>16 pt) is left unchanged
def _bs(size):
    return round(size * BODY * 2) / 2 if size <= 16 else size


def para(tf, text, size, color, bold=False, italic=False, font=SANS,
         align=PP_ALIGN.LEFT, after=4, first=False, line=None):
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    p.alignment = align
    if after is not None:
        p.space_after = Pt(after)
    p.space_before = Pt(0)
    if line is not None:
        p.line_spacing = line
    r = p.add_run()
    r.text = text
    f = r.font
    f.size = Pt(_bs(size))
    f.bold = bold
    f.italic = italic
    f.name = font
    f.color.rgb = color
    return p


def run(p, text, size, color, bold=False, italic=False, font=SANS):
    r = p.add_run()
    r.text = text
    f = r.font
    f.size = Pt(_bs(size))
    f.bold = bold
    f.italic = italic
    f.name = font
    f.color.rgb = color
    return r


def card(s, x, y, w, h, fill=PAPER2, line=LINE, line_w=1.0, radius=0.06):
    sh = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                            Inches(x), Inches(y), Inches(w), Inches(h))
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill
    if line is None:
        sh.line.fill.background()
    else:
        sh.line.color.rgb = line
        sh.line.width = Pt(line_w)
    sh.shadow.inherit = False
    try:
        sh.adjustments[0] = radius
    except Exception:
        pass
    return sh


def header(s, eyebrow, title, tsize=30):
    tf = box(s, 0.85, 0.62, 11.6, 0.35)
    para(tf, eyebrow.upper(), 11, GREEN, bold=True, font=SANS, after=0, first=True)
    tf2 = box(s, 0.85, 1.0, 11.6, 1.0)
    para(tf2, title, tsize, INK, bold=True, font=SERIF, after=0, first=True, line=1.04)


def bullets(s, x, y, w, h, items, size=12.5, color=INK_SOFT, gap=7,
            dot=GREEN, font=SANS):
    tf = box(s, x, y, w, h)
    for i, it in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(gap)
        p.space_before = Pt(0)
        p.line_spacing = 1.0
        r = p.add_run()
        r.text = "•  "
        r.font.size = Pt(_bs(size))
        r.font.name = font
        r.font.color.rgb = dot
        r.font.bold = True
        r2 = p.add_run()
        r2.text = it
        r2.font.size = Pt(_bs(size))
        r2.font.name = font
        r2.font.color.rgb = color


def stat(s, x, y, w, h, value, label, unit=None):
    card(s, x, y, w, h)
    tf = box(s, x + 0.2, y + 0.14, w - 0.4, h - 0.28, anchor=MSO_ANCHOR.MIDDLE)
    p = para(tf, value, 25, GREEN, bold=True, font=SERIF, after=2, first=True)
    if unit:
        run(p, " " + unit, 12, INK_SOFT, bold=True, font=SANS)
    para(tf, label, 10.5, INK_SOFT, font=SANS, after=0, line=1.0)


def tag(s, x, y, w, h, text, fill=TINT, color=GREEN):
    sh = card(s, x, y, w, h, fill=fill, line=None, radius=0.5)
    tf = box(s, x, y, w, h, anchor=MSO_ANCHOR.MIDDLE)
    para(tf, text.upper(), 9, color, bold=True, font=SANS, after=0,
         first=True, align=PP_ALIGN.CENTER)
    return sh


# ============================================================ 01 COVER
s = new_slide("verde-title.png")
tf = box(s, 0.95, 1.35, 11.0, 0.7)
para(tf, "UNIVERSIDAD NACIONAL DE COLOMBIA", 12, GOLD, bold=True, after=3, first=True)
para(tf, "MASTER'S THESIS PROPOSAL  ·  DEFENSE", 11.5, WHITE_SOFT, bold=False, after=0)
tf = box(s, 0.95, 2.35, 10.8, 1.7)
para(tf, "GeoAI-based Digital Twin for Dynamic Modeling of Coastal Ecosystems",
     38, WHITE, bold=True, font=SERIF, after=0, first=True, line=1.04)
tf = box(s, 0.95, 3.95, 8.6, 0.8)
para(tf, "Dynamic monitoring of the Ciénaga Grande de Santa Marta under persistent "
     "tropical cloud cover and hydroclimatic variability.", 15, WHITE_SOFT,
     font=SANS, after=0, first=True, line=1.12)
tf = box(s, 0.95, 4.85, 11.0, 1.1)
para(tf, "Lina Quintero Fonseca", 18, WHITE, bold=True, font=SERIF, after=3, first=True)
para(tf, "Master's in Geomatics  ·  Facultad de Ciencias Agrarias  ·  Sede Bogotá",
     11.5, WHITE_SOFT, after=2)
para(tf, "Advisor: Iván Lizarazo     ·     June 16, 2026", 11.5, WHITE_SOFT, after=0)
notes(s, "Buenos días. Agradezco al jurado y a mi director por su tiempo. Soy Lina "
       "Quintero Fonseca y presento mi propuesta de tesis de la Maestría en Geomática, "
       "titulada Gemelo Digital basado en GeoAI para el modelamiento dinámico de ecosistemas costeros. En "
       "esta propuesta, el marco conceptual del Gemelo Digital se desarrolla y evalúa "
       "tomando como caso de estudio la Ciénaga Grande de Santa Marta, con énfasis en el "
       "monitoreo dinámico del manglar, su condición y la inundación asociada. El título "
       "plantea un marco transferible; la aplicación empírica está delimitada a la CGSM.")

# ============================================================ 02 CONTEXT
s = new_slide("verde-content.png")
header(s, "Context · Why it matters", "The Ciénaga Grande de Santa Marta")
tf = box(s, 0.85, 2.0, 5.95, 1.05)
para(tf, "A strategic coastal lagoon-delta system in the Colombian Caribbean, with "
     "documented mangrove dieback, restoration efforts and persistent monitoring "
     "challenges.", 13, INK_SOFT, first=True, line=1.12)
stat(s, 0.85, 3.25, 2.85, 1.4, "1,286", "Core lagoon-mangrove system (km²)", unit=None)
stat(s, 3.85, 3.25, 2.95, 1.4, ">24,000", "Historical mangrove dieback (ha)")
stat(s, 0.85, 4.78, 2.85, 1.4, "~300,000", "People linked to ecosystem services")
stat(s, 3.85, 4.78, 2.95, 1.4, "5–12 d", "Satellite revisit potential")
# right gap panel
card(s, 7.05, 2.0, 5.45, 4.18, fill=PAPER2, line=LINE)
tf = box(s, 7.3, 2.25, 4.95, 0.4)
para(tf, "The monitoring gap", 15, INK, bold=True, font=SANS, after=0, first=True)
bullets(s, 7.3, 2.85, 4.95, 2.0, [
    "Fixed stations provide valuable but spatially limited information.",
    "Persistent cloud cover limits optical monitoring.",
    "SAR improves continuity, but vegetation interpretation remains complex.",
], size=12, gap=8, dot=GOLD)
card(s, 7.3, 5.18, 4.95, 0.82, fill=TINT, line=None)
tf = box(s, 7.5, 5.32, 4.55, 0.6, anchor=MSO_ANCHOR.MIDDLE)
p = para(tf, "Recent canopy loss in the ", 11.5, INK, after=0, first=True, line=1.05)
run(p, "Aguas Negras", 11.5, GREEN, bold=True)
run(p, " sector shows the need for spatially explicit monitoring.", 11.5, INK)
notes(s, "La CGSM es un sistema estratégico del Caribe colombiano. Su importancia no está "
       "solo en la extensión del manglar, sino en los servicios ecosistémicos que sostiene "
       "y en su historia de degradación y recuperación. El problema de fondo es que hay "
       "datos satelitales frecuentes, pero esa información no siempre se convierte en "
       "monitoreo espacialmente explícito y útil para la gestión. Las estaciones son "
       "fundamentales, pero no capturan toda la heterogeneidad espacial del sistema.")

# ============================================================ 03 STUDY AREA
s = new_slide("verde-content.png")
header(s, "Study area · Data readiness", "A bounded AOI with available multi-sensor data")
items = [
    ("AOI — Core CGSM study area",
     "1,286 km² lagoon-mangrove polygon within the broader Ramsar system."),
    ("Imagery availability checked in GEE",
     "Sentinel-2, Sentinel-1 C-band, ALOS-2 L-band, Landsat 8/9 and ancillary layers for 2015–2025."),
    ("Reference layers and metric georeferencing",
     "GMW v3.0, INVEMAR stations, CARICOMP and SRTM; metric areas in EPSG:9377 (MAGNA-SIRGAS)."),
]
yy = 2.2
for ti, de in items:
    card(s, 0.85, yy, 5.0, 1.18)
    tf = box(s, 1.08, yy + 0.16, 4.55, 0.9, anchor=MSO_ANCHOR.MIDDLE)
    para(tf, ti, 12.5, INK, bold=True, font=SANS, after=2, first=True)
    para(tf, de, 10.5, INK_SOFT, after=0, line=1.05)
    yy += 1.34
s.shapes.add_picture(os.path.join(FIG, "study-area.png"), Inches(6.2), Inches(2.05),
                     height=Inches(4.4))
notes(s, "Para evitar una propuesta abstracta, el área de estudio se delimitó desde el "
       "inicio. El AOI es el núcleo lagunar-manglar de la Ciénaga Grande de Santa Marta, "
       "aproximadamente 1.286 kilómetros cuadrados. Uso el término sensu stricto para "
       "aclarar que la tesis no analiza todo el sitio Ramsar ni todo el delta del Magdalena, "
       "sino ese núcleo definido como área de estudio. Además, se revisó la disponibilidad "
       "de imágenes y capas de referencia en Google Earth Engine para asegurar que la "
       "metodología sea viable desde la fase inicial. Las áreas métricas se trabajarán en "
       "EPSG:9377, que corresponde al sistema oficial colombiano.")

# ============================================================ 04 GAP
s = new_slide("verde-content.png")
header(s, "Why this research", "Three gaps constrain operational monitoring")
gaps = [
    ("1", "Data infrastructure",
     "No analysis-ready optical-SAR datacube has been implemented for CGSM monitoring."),
    ("2", "GeoAI benchmarking",
     "SAM-based segmentation has not been systematically evaluated for CGSM mangrove "
     "and flood mapping."),
    ("3", "Operational integration",
     "Monitoring products remain weakly connected to automated, reproducible Digital "
     "Twin workflows."),
]
xx = 0.85
for n, ti, de in gaps:
    card(s, xx, 2.15, 3.78, 2.85)
    tf = box(s, xx + 0.28, 2.4, 3.25, 2.4)
    para(tf, n.zfill(2), 26, GREEN, bold=True, font=SERIF, after=4, first=True)
    para(tf, ti, 14, INK, bold=True, font=SANS, after=6, line=1.0)
    para(tf, de, 11, INK_SOFT, after=0, line=1.08)
    xx += 3.95
card(s, 0.85, 5.35, 11.65, 0.95, fill=TINT, line=None)
tf = box(s, 1.1, 5.5, 11.15, 0.7, anchor=MSO_ANCHOR.MIDDLE)
p = para(tf, "Opportunity: ", 12.5, GREEN, bold=True, after=0, first=True, line=1.05)
run(p, "connect satellite archives, GeoAI models and monitoring outputs in one "
       "evaluable framework.", 12.5, INK)
notes(s, "La brecha no está en la ausencia de datos. La brecha está en la integración. Hay "
       "datos ópticos, radar y referencias institucionales, pero falta una arquitectura que "
       "los organice como datacube, evalúe modelos GeoAI y entregue productos reproducibles "
       "para monitoreo. La oportunidad de esta investigación está en conectar esos "
       "componentes en un solo flujo evaluable.")

# ============================================================ 05 OBJECTIVES
s = new_slide("verde-content.png")
header(s, "Research question · Objectives", "One question, three objectives")
card(s, 0.85, 2.05, 11.65, 1.45, fill=GREEN, line=None)
tf = box(s, 1.15, 2.2, 11.05, 1.18, anchor=MSO_ANCHOR.MIDDLE)
para(tf, "How can optical and radar Earth observation be integrated into a GeoAI-based "
     "Digital Twin to monitor mangrove extent, condition and flooding dynamics in the CGSM?",
     16, WHITE, font=SERIF, after=0, first=True, line=1.14)
objs = [
    ("Objective 1", "Design the datacube",
     "Integrate multi-sensor optical and radar data into analysis-ready composites."),
    ("Objective 2", "Benchmark GeoAI",
     "Compare SAM-based segmentation against a Random Forest baseline."),
    ("Objective 3", "Build the Digital Twin",
     "Integrate outputs into reusable, automated and reproducible workflows."),
]
xx = 0.85
for tg, ti, de in objs:
    card(s, xx, 3.75, 3.78, 2.55)
    tf = box(s, xx + 0.26, 3.98, 3.28, 2.15)
    para(tf, tg.upper(), 9.5, GREEN, bold=True, after=5, first=True)
    para(tf, ti, 14, INK, bold=True, font=SANS, after=6, line=1.0)
    para(tf, de, 11, INK_SOFT, after=0, line=1.08)
    xx += 3.95
notes(s, "La pregunta se concentra en cómo integrar datos ópticos y radar dentro de un "
       "Gemelo Digital basado en GeoAI para monitorear tres variables: extensión, condición "
       "e inundación. La nubosidad persistente y la variabilidad hidroclimática justifican "
       "el uso de fusión óptico-radar. Los tres objetivos responden directamente a esa "
       "pregunta: primero organizo los datos, luego evalúo los modelos y finalmente integro "
       "los productos en un prototipo reproducible.")

# ============================================================ 06 COHERENCE
s = new_slide("verde-content.png")
header(s, "Verification checklist", "One coherent thread")
rows = [
    ("Element", "What it answers"),
    ("Title", "GeoAI Digital Twin for dynamic coastal ecosystem modeling"),
    ("Case study", "CGSM mangrove and flooding dynamics"),
    ("Problem", "Fragmented optical-radar monitoring"),
    ("Question", "Integration of EO data, GeoAI and Digital Twin"),
    ("Objective 1", "Datacube"),
    ("Objective 2", "Model benchmark"),
    ("Objective 3", "Reproducible prototype"),
]
gt = s.shapes.add_table(8, 2, Inches(0.85), Inches(2.05), Inches(11.65), Inches(3.7)).table
gt.first_row = False
gt.horz_banding = False
gt.columns[0].width = Inches(3.2)
gt.columns[1].width = Inches(8.45)
for r in range(8):
    for c in range(2):
        cl = gt.cell(r, c)
        cl.margin_left = Inches(0.14)
        cl.margin_right = Inches(0.1)
        cl.margin_top = Inches(0.04)
        cl.margin_bottom = Inches(0.04)
        cl.vertical_anchor = MSO_ANCHOR.MIDDLE
        cl.fill.solid()
        if r == 0:
            cl.fill.fore_color.rgb = GREEN
        else:
            cl.fill.fore_color.rgb = WHITE if r % 2 == 1 else PAPER2
        pc = cl.text_frame.paragraphs[0]
        rc = pc.add_run()
        rc.text = rows[r][c]
        fc = rc.font
        fc.name = SANS
        fc.size = Pt(11.5)
        if r == 0:
            fc.bold = True
            fc.color.rgb = WHITE
        else:
            fc.bold = (c == 0)
            fc.color.rgb = GREEN if c == 0 else INK
tf = box(s, 0.85, 5.95, 11.65, 0.4)
para(tf, "Each objective has a method, an output and an evaluation criterion.",
     11.5, INK_SOFT, first=True)
notes(s, "Esta diapositiva responde a una preocupación típica del jurado: la coherencia "
       "interna. La propuesta no presenta técnicas sueltas. Cada objetivo tiene una "
       "pregunta, un método, un resultado y una forma de evaluación. El objetivo 1 produce "
       "la infraestructura de datos; el objetivo 2 produce los mapas y la comparación de "
       "modelos; y el objetivo 3 integra esos productos en un prototipo reproducible.")

# ============================================================ 07 FRAMEWORK
s = new_slide("verde-content.png")
header(s, "Methodology · Operational framework", "From satellite observations to monitoring outputs")
s.shapes.add_picture(os.path.join(FIG, "framework.png"), Inches(1.6), Inches(1.95),
                     width=Inches(10.13))
notes(s, "La metodología se organiza en tres capas. La primera es la capa de observación: "
       "sensores ópticos, radar y datos de referencia. La segunda es la capa analítica: "
       "datacube, Random Forest, modelos SAM-based y detección de inundación. La tercera es "
       "la capa operativa: flujo automatizado, tablero y repositorio. La validación conecta "
       "las tres capas con el Global Mangrove Watch, el INVEMAR y CARICOMP.")

# ============================================================ 08 DATA & ROUTE
s = new_slide("verde-content.png")
header(s, "Methods · Objectives 1 & 2", "Datacube + prioritized modeling route")
card(s, 0.85, 2.1, 5.7, 4.05)
tf = box(s, 1.1, 2.32, 5.2, 0.4)
para(tf, "Analysis-ready datacube", 14, INK, bold=True, after=0, first=True)
bullets(s, 1.1, 2.95, 5.25, 3.0, [
    "10 m working grid",
    "Dry-season optical composites",
    "Seasonal SAR composites",
    "NDVI · NDWI · CMRI · NDBaI",
], size=12.5, gap=8)
# datacube schematic — fills the lower half of the left panel
tf = box(s, 1.1, 4.25, 5.3, 0.3)
para(tf, "THE DATACUBE · LAYERS CO-REGISTERED AT 10 m · 2015–2025", 8.5, INK_SOFT,
     bold=True, after=0, first=True)
_layers = [
    ("Sentinel-2 · optical", GREEN),
    ("Sentinel-1 · C-band SAR", RGBColor(0x3b, 0x82, 0xc4)),
    ("ALOS-2 · L-band SAR", GOLD),
    ("Indices: NDVI · NDWI · CMRI", RGBColor(0x14, 0xb0, 0x86)),
]
_ly = 4.66
for _name, _col in _layers:
    card(s, 1.15, _ly, 5.1, 0.34, fill=PAPER2, line=LINE, radius=0.16)
    _a = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(1.34), Inches(_ly + 0.1),
                            Inches(0.15), Inches(0.15))
    _a.fill.solid(); _a.fill.fore_color.rgb = _col
    _a.line.fill.background(); _a.shadow.inherit = False
    _t = box(s, 1.66, _ly, 4.4, 0.34, anchor=MSO_ANCHOR.MIDDLE)
    para(_t, _name, 10.5, INK, after=0, first=True)
    _ly += 0.38
card(s, 6.8, 2.1, 5.7, 4.05, fill=GREEN, line=None)
tf = box(s, 7.05, 2.32, 5.2, 0.4)
para(tf, "Chosen route: Random Forest", 14, WHITE, bold=True, after=0, first=True)
tf = box(s, 7.05, 2.92, 5.2, 0.7)
para(tf, "Primary classifier on the fused datacube — 200 trees, smileRandomForest in GEE.",
     12.5, WHITE_SOFT, after=0, first=True, line=1.12)
bullets(s, 7.05, 3.72, 5.25, 2.2, [
    "Widely used in tropical mangrove mapping",
    "Interpretable; transparent variable importance (Gini)",
    "Robust with limited, heterogeneous samples",
    "GEE-native — runs at cloud scale, no local setup",
], size=12, gap=13, dot=GOLD, color=WHITE_SOFT)
tf = box(s, 7.05, 5.62, 5.25, 0.5)
p = para(tf, "Benchmarked against ", 11.5, WHITE_SOFT, after=0, first=True, line=1.08)
run(p, "SamGeo / MW-SAM", 11.5, GOLD, bold=True)
run(p, " to assess the added value of radar and segmentation.", 11.5, WHITE_SOFT)
notes(s, "Sobre los métodos, priorizo una ruta concreta. Primero construyo un datacube "
       "listo para análisis, con composiciones ópticas y SAR, índices espectrales y "
       "co-registro espacial. La ruta principal de clasificación es Random Forest porque es "
       "interpretable, corre en Earth Engine y funciona bien con datos heterogéneos. SamGeo "
       "y MW-SAM se evalúan como comparación, para valorar el aporte del radar y de la "
       "segmentación, no como promesa de superioridad.")

# ============================================================ 09 METRICS
s = new_slide("verde-content.png")
header(s, "Methods · Mapping & evaluation criteria", "Classification, floods & metrics")
card(s, 0.85, 2.1, 5.7, 4.05)
tf = box(s, 1.1, 2.32, 5.2, 0.4)
para(tf, "Classes & flood mapping", 14, INK, bold=True, after=0, first=True)
tag(s, 1.1, 2.95, 1.15, 0.34, "Intact")
tag(s, 2.35, 2.95, 1.45, 0.34, "Degraded")
tag(s, 3.9, 2.95, 1.75, 0.34, "Non-mangrove")
bullets(s, 1.1, 3.6, 5.25, 2.4, [
    "C-band for open / degraded areas",
    "L-band support for closed canopy",
    "Layers merged by canopy-closure stratification",
], size=12, gap=10)
card(s, 6.8, 2.1, 5.7, 4.05)
tf = box(s, 7.05, 2.32, 5.2, 0.4)
para(tf, "Evaluation criteria", 14, INK, bold=True, after=0, first=True)
mets = [
    ("Classification (RF)", "Overall Accuracy · per-class Precision / Recall / F1"),
    ("Segmentation (SAM)", "Intersection-over-Union (IoU)"),
    ("Continuous (canopy height)", "RMSE · R²"),
]
yy = 2.95
for ti, de in mets:
    card(s, 7.05, yy, 5.2, 0.82, fill=PAPER2, line=LINE)
    tf = box(s, 7.28, yy + 0.12, 4.75, 0.6, anchor=MSO_ANCHOR.MIDDLE)
    para(tf, ti, 11.5, INK, bold=True, after=2, first=True, line=1.0)
    para(tf, de, 10.5, INK_SOFT, after=0, line=1.0)
    yy += 0.92
card(s, 7.05, yy, 5.2, 0.62, fill=TINT, line=None)
tf = box(s, 7.28, yy + 0.07, 4.75, 0.48, anchor=MSO_ANCHOR.MIDDLE)
p = para(tf, "Kappa only as a secondary metric. Validated against ", 10.5, INK, after=0, first=True, line=1.0)
run(p, "GMW v3.0", 10.5, GREEN, bold=True)
run(p, ", INVEMAR and CARICOMP.", 10.5, INK)
notes(s, "Cada producto tiene una métrica. Para clasificación no me quedo solo con exactitud "
       "global, porque puede ocultar errores en clases difíciles. Por eso uso precisión, "
       "recall y F1 por clase. Para segmentación uso IoU, y para variables continuas usaría "
       "RMSE y R cuadrado. Kappa puede reportarse como métrica secundaria, pero no será la "
       "base de interpretación.")

# ============================================================ 10 WORKED EXAMPLE
s = new_slide("verde-content.png")
header(s, "Worked example · Rinconada sector", "From satellite pixels to a monitoring layer")
steps = [
    ("1", "Ingest", "S2 dry-season, S1 wet-season and ALOS-2 L-band composites for the Rinconada sector."),
    ("2", "Classify", "RF labels each pixel intact / degraded / non-mangrove; SAM benchmarked."),
    ("3", "Detect floods", "L-band support under closed canopy; C-band in open / degraded areas."),
    ("4", "Support monitoring", "Flood frequency by sector; comparison against the 2015–2025 baseline."),
]
xx = 0.85
cw = 2.75
for n, ti, de in steps:
    card(s, xx, 2.95, cw, 2.55)
    tf = box(s, xx + 0.26, 3.2, cw - 0.5, 2.2)
    para(tf, n, 24, GREEN, bold=True, font=SERIF, after=5, first=True)
    para(tf, ti, 13.5, INK, bold=True, font=SANS, after=6, line=1.0)
    para(tf, de, 10.5, INK_SOFT, after=0, line=1.12)
    xx += cw + 0.3
notes(s, "Este ejemplo evita que la metodología suene abstracta. Muestra el flujo completo: "
       "primero se ingieren los datos, luego se clasifica la cobertura, después se detecta "
       "inundación según el tipo de dosel y finalmente el prototipo muestra frecuencia de "
       "inundación por sector. El producto no reemplaza una decisión institucional, pero sí "
       "genera una capa de apoyo para priorizar monitoreo.")

# ============================================================ 11 RESULTS
s = new_slide("verde-content.png")
header(s, "Expected results · 24 months", "Outputs, prototype and plan")
outs = [
    ("Obj. 1 — Datacube", "Multi-sensor analysis-ready datacube, 2015–2025, 10 m, on GEE."),
    ("Obj. 2 — GeoAI maps", "Validated mangrove + flood maps; RF-vs-SAM comparison."),
    ("Obj. 3 — Digital Twin", "Interactive prototype + open pipeline, article & technical report."),
]
xx = 0.85
for ti, de in outs:
    card(s, xx, 2.0, 3.78, 1.15)
    tf = box(s, xx + 0.24, 2.16, 3.34, 0.85)
    para(tf, ti, 12.5, INK, bold=True, after=3, first=True)
    para(tf, de, 10.5, INK_SOFT, after=0, line=1.05)
    xx += 3.95
# timeline table (left)
acts = [
    ("Datacube design & construction", [1, 1, 0, 0]),
    ("RF + SAM/MW-SAM experiments", [0, 1, 1, 0]),
    ("Flood dynamics & change detection", [0, 1, 1, 0]),
    ("Digital Twin & dashboard", [0, 0, 1, 1]),
    ("Writing · article · defense", [0, 0, 1, 1]),
]
gt = s.shapes.add_table(6, 5, Inches(0.85), Inches(3.4), Inches(11.65), Inches(2.7)).table
gt.first_row = False
gt.horz_banding = False
for i, cw in enumerate([6.05, 1.4, 1.4, 1.4, 1.4]):
    gt.columns[i].width = Inches(cw)
hdr = ["Activity", "S1", "S2", "S3", "S4"]
for c in range(5):
    cl = gt.cell(0, c)
    cl.fill.solid()
    cl.fill.fore_color.rgb = PAPER2
    cl.vertical_anchor = MSO_ANCHOR.MIDDLE
    cl.margin_top = Inches(0.03)
    cl.margin_bottom = Inches(0.03)
    p = cl.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT if c == 0 else PP_ALIGN.CENTER
    rr = p.add_run()
    rr.text = hdr[c]
    rr.font.size = Pt(10)
    rr.font.bold = True
    rr.font.name = SANS
    rr.font.color.rgb = INK_SOFT
for r, (lab, cells) in enumerate(acts, start=1):
    cl = gt.cell(r, 0)
    cl.fill.solid()
    cl.fill.fore_color.rgb = WHITE
    cl.vertical_anchor = MSO_ANCHOR.MIDDLE
    cl.margin_left = Inches(0.08)
    p = cl.text_frame.paragraphs[0]
    rr = p.add_run()
    rr.text = lab
    rr.font.size = Pt(10.5)
    rr.font.name = SANS
    rr.font.color.rgb = INK
    for j, on in enumerate(cells):
        cc = gt.cell(r, j + 1)
        cc.fill.solid()
        cc.fill.fore_color.rgb = GREEN if on else WHITE
        cc.text_frame.paragraphs[0].add_run().text = ""
notes(s, "Los resultados esperados se organizan por objetivo. El primero entrega el "
       "datacube; el segundo, los mapas validados y el benchmark; el tercero, el prototipo "
       "interactivo y el repositorio. El cronograma se distribuye en cuatro semestres: "
       "construcción del datacube, experimentos de modelos, análisis de inundación y "
       "cambio, integración del prototipo y escritura. La latencia objetivo y la capacidad "
       "computacional quedan documentadas como referencia: el procesamiento se apoya en "
       "Earth Engine, con requerimientos limitados de infraestructura local.")

# ============================================================ 12 CLOSING
s = new_slide("verde-title.png")
tf = box(s, 0.95, 1.2, 11.0, 0.4)
para(tf, "TAKE-HOME MESSAGES", 12, GOLD, bold=True, after=0, first=True)
msgs = [
    ("1", "The problem is operational", "Data exist, but monitoring remains fragmented."),
    ("2", "The method is evaluable", "Each output has a target variable, reference data and a metric."),
    ("3", "The contribution is transferable", "The workflow can be adapted to other tropical coastal wetlands."),
]
yy = 2.05
for n, t, d in msgs:
    nb = box(s, 0.95, yy, 0.7, 0.9)
    para(nb, n, 26, GOLD, bold=True, font=SERIF, after=0, first=True)
    tb = box(s, 1.75, yy + 0.04, 10.5, 0.9)
    para(tb, t, 17, WHITE, bold=True, font=SERIF, after=2, first=True)
    para(tb, d, 12.5, WHITE_SOFT, after=0, line=1.05)
    yy += 1.0
tf = box(s, 0.95, 5.0, 11.0, 0.95)
para(tf, "Thank you", 24, WHITE, bold=True, font=SERIF, after=3, first=True)
para(tf, "Questions and comments", 12, WHITE_SOFT, after=0)
notes(s, "Cierro con tres ideas. Primero, el problema no es la ausencia de datos, sino su "
       "integración operativa. Segundo, la metodología es evaluable porque cada producto "
       "tiene datos, referencia y métrica. Tercero, la contribución es transferible como "
       "arquitectura, no como un modelo que se copia sin recalibrar. Muchas gracias, quedo "
       "atenta a sus preguntas y comentarios.")

# ============================================================ 13 BACKUP
s = new_slide("verde-content.png")
header(s, "Backup · Feasibility", "Timeline & budget")
card(s, 0.85, 2.1, 5.7, 4.05)
tf = box(s, 1.1, 2.32, 5.2, 0.4)
para(tf, "Timeline", 14, INK, bold=True, after=0, first=True)
tf = box(s, 1.1, 2.85, 5.2, 0.5)
p = para(tf, "24", 26, GREEN, bold=True, font=SERIF, after=0, first=True)
run(p, "  months · 4 semesters", 12, INK_SOFT)
phases = [
    ("Phase 1", "Datacube construction", "S1–S2"),
    ("Phase 2", "GeoAI benchmark & flood mapping", "S2–S3"),
    ("Phase 3", "Digital Twin · writing · defense", "S3–S4"),
]
yy = 3.7
for pn, pt, ps in phases:
    tf = box(s, 1.1, yy, 4.1, 0.55, anchor=MSO_ANCHOR.MIDDLE)
    p = para(tf, pn.upper() + "   ", 9.5, GREEN, bold=True, after=0, first=True, line=1.0)
    run(p, pt, 11.5, INK)
    tg = box(s, 5.25, yy, 1.05, 0.55, anchor=MSO_ANCHOR.MIDDLE)
    para(tg, ps, 10, INK_SOFT, bold=True, after=0, first=True, align=PP_ALIGN.RIGHT)
    yy += 0.72
# budget panel
card(s, 6.8, 2.1, 5.7, 4.05)
tf = box(s, 7.05, 2.32, 5.2, 0.4)
para(tf, "Budget estimate", 14, INK, bold=True, after=0, first=True)
tf = box(s, 7.05, 2.85, 5.2, 0.5)
p = para(tf, "~15–18 M", 25, GREEN, bold=True, font=SERIF, after=0, first=True)
run(p, "  COP · estimated operating budget", 10.5, INK_SOFT)
tf = box(s, 7.05, 3.52, 5.2, 0.3)
para(tf, "MAIN COST DRIVERS", 9.5, INK_SOFT, bold=True, after=0, first=True)
bullets(s, 7.05, 3.96, 5.3, 1.7, [
    "Technical support",
    "Storage & backup",
    "Validation & institutional meetings",
    "Publication / dissemination support",
], size=11.5, gap=9)
tf = box(s, 7.05, 5.62, 5.3, 0.5)
p = para(tf, "Academic supervision and institutional infrastructure are treated as ",
         10, INK_SOFT, after=0, first=True, line=1.05)
run(p, "in-kind support.", 10, INK, bold=True)
# bottom note
card(s, 0.85, 6.32, 11.65, 0.72, fill=TINT, line=None)
tf = box(s, 1.1, 6.42, 11.15, 0.55, anchor=MSO_ANCHOR.MIDDLE)
p = para(tf, "Mostly based on open data and institutional infrastructure ", 11, GREEN,
         bold=True, after=0, first=True, line=1.0)
run(p, "— open satellite archives (Sentinel, Landsat, ALOS-2, GMW), Google Earth "
       "Engine, and UNAL computing. No commercial imagery is required.", 11, INK)
notes(s, "Esta diapositiva es de respaldo. El proyecto es viable porque se basa en datos "
       "abiertos, Google Earth Engine, software abierto e infraestructura institucional. El "
       "presupuesto operativo se concentra en apoyo técnico puntual, almacenamiento, "
       "validación o socialización institucional y divulgación. No requiere compra de "
       "imágenes comerciales.")

prs.save(OUT)
print("SAVED", OUT, "slides:", len(prs.slides._sldIdLst))
